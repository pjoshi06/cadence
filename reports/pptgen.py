"""Builds DSR/WSR/MSR PowerPoint decks from ticket data.

Works in two modes:
  * built-in default template — a clean deck styled in code
  * uploaded client template  — slides are appended using the client's
    slide master/layouts so their branding carries through
"""
import io
from datetime import datetime, time, timedelta

from django.utils import timezone
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

from .models import Ticket

# Palette (slate / indigo — matches the portal UI)
INK = RGBColor(0x1E, 0x29, 0x3B)        # slate-800
MUTED = RGBColor(0x64, 0x74, 0x8B)      # slate-500
ACCENT = RGBColor(0x4F, 0x46, 0xE5)     # indigo-600
ACCENT_DARK = RGBColor(0x31, 0x2E, 0x81)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)      # slate-100
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOOD = RGBColor(0x05, 0x96, 0x69)       # emerald-600
BAD = RGBColor(0xDC, 0x26, 0x26)        # red-600
WARN = RGBColor(0xD9, 0x77, 0x06)      # amber-600

PRIORITY_COLORS = {"P1": BAD, "P2": WARN, "P3": ACCENT, "P4": MUTED}


# ---------------------------------------------------------------- metrics

def compute_metrics(period_start, period_end):
    tz = timezone.get_current_timezone()
    start_dt = timezone.make_aware(datetime.combine(period_start, time.min), tz)
    end_dt = timezone.make_aware(datetime.combine(period_end, time.max), tz)

    opened = Ticket.objects.filter(opened_at__range=(start_dt, end_dt))
    resolved = Ticket.objects.filter(resolved_at__range=(start_dt, end_dt))
    backlog = Ticket.objects.filter(opened_at__lte=end_dt).exclude(
        state__in=Ticket.OPEN_EXCLUDED_STATES
    )

    sla_base = resolved.exclude(sla_met__isnull=True)
    sla_met = sla_base.filter(sla_met=True).count()
    sla_pct = round(sla_met / sla_base.count() * 100, 1) if sla_base.count() else None

    prio = {p: backlog.filter(priority=p).count() for p in ("P1", "P2", "P3", "P4")}

    ageing = {"0–3 days": 0, "4–7 days": 0, "8–15 days": 0, "> 15 days": 0}
    for t in backlog:
        d = t.age_days
        if d <= 3:
            ageing["0–3 days"] += 1
        elif d <= 7:
            ageing["4–7 days"] += 1
        elif d <= 15:
            ageing["8–15 days"] += 1
        else:
            ageing["> 15 days"] += 1

    focus = list(
        backlog.filter(priority__in=("P1", "P2")).order_by("priority", "opened_at")[:10]
    )

    return {
        "opened": opened.count(),
        "resolved": resolved.count(),
        "backlog": backlog.count(),
        "sla_pct": sla_pct,
        "priority": prio,
        "ageing": ageing,
        "focus": focus,
        "resolved_list": list(resolved.order_by("-resolved_at")[:10]),
    }


# ---------------------------------------------------------------- helpers

def _blank_layout(prs):
    """Pick the emptiest layout so client branding shows but placeholders don't."""
    best = prs.slide_layouts[len(prs.slide_layouts) - 1]
    best_count = 99
    for layout in prs.slide_layouts:
        n = len(layout.placeholders)
        if n < best_count:
            best, best_count = layout, n
    return best


def _add_slide(prs):
    slide = prs.slides.add_slide(_blank_layout(prs))
    # drop any inherited placeholders so slides stay clean
    for shape in list(slide.placeholders):
        shape.element.getparent().remove(shape.element)
    return slide


def _textbox(slide, left, top, width, height, text, size, color=INK,
             bold=False, align=PP_ALIGN.LEFT, font="Calibri"):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = font
    return box


def _rect(slide, left, top, width, height, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _header(slide, prs, title, subtitle, branded):
    w = prs.slide_width
    if not branded:
        _rect(slide, 0, 0, w, Inches(0.16), ACCENT)
    _textbox(slide, Inches(0.5), Inches(0.35), w - Inches(1.0), Inches(0.5),
             title, 24, INK, bold=True)
    _textbox(slide, Inches(0.5), Inches(0.82), w - Inches(1.0), Inches(0.35),
             subtitle, 12, MUTED)


def _kpi_card(slide, left, top, width, height, label, value, color=INK, note=""):
    _rect(slide, left, top, width, height, LIGHT)
    _textbox(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4),
             Inches(0.3), label.upper(), 10, MUTED, bold=True)
    _textbox(slide, left + Inches(0.2), top + Inches(0.45), width - Inches(0.4),
             Inches(0.6), value, 30, color, bold=True)
    if note:
        _textbox(slide, left + Inches(0.2), top + height - Inches(0.4),
                 width - Inches(0.4), Inches(0.3), note, 9, MUTED)


def _style_table(table, header_fill=ACCENT_DARK):
    for cell in table.rows[0].cells:
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE
                r.font.bold = True
                r.font.size = Pt(11)
    for row in list(table.rows)[1:]:
        for cell in row.cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10)
                    r.font.color.rgb = INK
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def _add_table(slide, rows, cols, left, top, width, height):
    return slide.shapes.add_table(rows, cols, left, top, width, height).table


def _set_cell(table, r, c, text, bold=False, color=None):
    cell = table.cell(r, c)
    cell.text = str(text)
    for p in cell.text_frame.paragraphs:
        for run in p.runs:
            run.font.bold = bold
            if color:
                run.font.color.rgb = color


# ---------------------------------------------------------------- slides

def _title_slide(prs, kind_label, period_label, team_name, branded):
    slide = _add_slide(prs)
    w, h = prs.slide_width, prs.slide_height
    if not branded:
        _rect(slide, 0, 0, w, h, ACCENT_DARK)
        _rect(slide, 0, int(h * 0.62), w, Inches(0.06), ACCENT)
        text_color, sub_color = WHITE, RGBColor(0xC7, 0xD2, 0xFE)
    else:
        text_color, sub_color = INK, MUTED
    _textbox(slide, Inches(0.8), int(h * 0.30), w - Inches(1.6), Inches(1.0),
             kind_label, 40, text_color, bold=True)
    _textbox(slide, Inches(0.8), int(h * 0.47), w - Inches(1.6), Inches(0.5),
             f"{team_name}  |  Reporting period: {period_label}", 16, sub_color)
    _textbox(slide, Inches(0.8), int(h * 0.68), w - Inches(1.6), Inches(0.4),
             f"Generated by Cadence on {timezone.localdate().strftime('%d %b %Y')}",
             11, sub_color)


def _summary_slide(prs, m, period_label, branded):
    slide = _add_slide(prs)
    _header(slide, prs, "Executive Summary", period_label, branded)
    w = prs.slide_width
    margin, gap = Inches(0.5), Inches(0.25)
    card_w = int((w - margin * 2 - gap * 3) / 4)
    top, card_h = Inches(1.4), Inches(1.5)
    sla_text = f"{m['sla_pct']}%" if m["sla_pct"] is not None else "N/A"
    sla_color = GOOD if (m["sla_pct"] or 0) >= 95 else (WARN if (m["sla_pct"] or 0) >= 85 else BAD)
    cards = [
        ("Tickets Received", str(m["opened"]), INK),
        ("Tickets Resolved", str(m["resolved"]), GOOD),
        ("Open Backlog", str(m["backlog"]), WARN if m["backlog"] else GOOD),
        ("SLA Compliance", sla_text, sla_color if m["sla_pct"] is not None else MUTED),
    ]
    for i, (label, value, color) in enumerate(cards):
        _kpi_card(slide, margin + i * (card_w + gap), top, card_w, card_h, label, value, color)

    # priority breakdown bars
    _textbox(slide, margin, Inches(3.15), Inches(4), Inches(0.35),
             "OPEN BACKLOG BY PRIORITY", 11, MUTED, bold=True)
    total = sum(m["priority"].values()) or 1
    bar_top = Inches(3.55)
    bar_full = w - margin * 2 - Inches(2.2)
    for p in ("P1", "P2", "P3", "P4"):
        n = m["priority"][p]
        _textbox(slide, margin, bar_top, Inches(0.6), Inches(0.3), p, 11, INK, bold=True)
        track = _rect(slide, margin + Inches(0.7), bar_top + Inches(0.03),
                      bar_full, Inches(0.22), LIGHT)
        if n:
            _rect(slide, margin + Inches(0.7), bar_top + Inches(0.03),
                  max(int(bar_full * n / total), Inches(0.12)), Inches(0.22),
                  PRIORITY_COLORS[p])
        _textbox(slide, margin + Inches(0.8) + bar_full, bar_top, Inches(1.2),
                 Inches(0.3), str(n), 11, INK, bold=True)
        bar_top += Inches(0.42)


def _ageing_slide(prs, m, period_label, branded):
    slide = _add_slide(prs)
    _header(slide, prs, "Backlog Ageing & Focus Items", period_label, branded)
    margin = Inches(0.5)

    table = _add_table(slide, 5, 2, margin, Inches(1.4), Inches(3.6), Inches(2.0))
    _set_cell(table, 0, 0, "Ageing Bucket")
    _set_cell(table, 0, 1, "Open Tickets")
    for i, (bucket, n) in enumerate(m["ageing"].items(), start=1):
        _set_cell(table, i, 0, bucket)
        _set_cell(table, i, 1, n, bold=n > 0 and bucket == "> 15 days",
                  color=BAD if (n and bucket == "> 15 days") else None)
    _style_table(table)

    _textbox(slide, Inches(4.5), Inches(1.35), Inches(5), Inches(0.35),
             "TOP OPEN P1 / P2 TICKETS", 11, MUTED, bold=True)
    focus = m["focus"]
    if focus:
        rows = len(focus) + 1
        ft = _add_table(slide, rows, 4, Inches(4.5), Inches(1.7),
                        prs.slide_width - Inches(5.0), Inches(0.35) * rows)
        for c, htxt in enumerate(["Ticket", "Priority", "Age (d)", "Summary"]):
            _set_cell(ft, 0, c, htxt)
        for r, t in enumerate(focus, start=1):
            _set_cell(ft, r, 0, t.number)
            _set_cell(ft, r, 1, t.priority, bold=True,
                      color=PRIORITY_COLORS.get(t.priority))
            _set_cell(ft, r, 2, t.age_days)
            desc = t.short_description
            _set_cell(ft, r, 3, desc[:60] + ("…" if len(desc) > 60 else ""))
        _style_table(ft)
    else:
        _textbox(slide, Inches(4.5), Inches(1.8), Inches(4), Inches(0.4),
                 "No open P1/P2 tickets. ✓", 13, GOOD, bold=True)


def _resolved_slide(prs, m, period_label, branded):
    slide = _add_slide(prs)
    _header(slide, prs, "Tickets Resolved in Period", period_label, branded)
    margin = Inches(0.5)
    resolved = m["resolved_list"]
    if not resolved:
        _textbox(slide, margin, Inches(1.6), Inches(6), Inches(0.4),
                 "No tickets were resolved in this period.", 13, MUTED)
        return
    rows = len(resolved) + 1
    table = _add_table(slide, rows, 5, margin, Inches(1.4),
                       prs.slide_width - margin * 2, Inches(0.35) * rows)
    for c, htxt in enumerate(["Ticket", "Priority", "Summary", "Resolved By", "SLA"]):
        _set_cell(table, 0, c, htxt)
    for r, t in enumerate(resolved, start=1):
        _set_cell(table, r, 0, t.number)
        _set_cell(table, r, 1, t.priority)
        desc = t.short_description
        _set_cell(table, r, 2, desc[:55] + ("…" if len(desc) > 55 else ""))
        _set_cell(table, r, 3, t.assigned_to or "—")
        if t.sla_met is None:
            _set_cell(table, r, 4, "—")
        else:
            _set_cell(table, r, 4, "Met" if t.sla_met else "Missed",
                      bold=True, color=GOOD if t.sla_met else BAD)
    _style_table(table)


def _highlights_slide(prs, highlights, period_label, branded):
    slide = _add_slide(prs)
    _header(slide, prs, "Highlights & Manager Notes", period_label, branded)
    margin = Inches(0.5)
    _rect(slide, margin, Inches(1.4), prs.slide_width - margin * 2,
          prs.slide_height - Inches(2.1), LIGHT)
    box = slide.shapes.add_textbox(margin + Inches(0.3), Inches(1.7),
                                   prs.slide_width - margin * 2 - Inches(0.6),
                                   prs.slide_height - Inches(2.6))
    tf = box.text_frame
    tf.word_wrap = True
    lines = [ln.strip() for ln in (highlights or "").splitlines() if ln.strip()]
    if not lines:
        lines = ["(No notes provided — edit this slide before sending.)"]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = f"•  {line}"
        run.font.size = Pt(14)
        run.font.color.rgb = INK
        p.space_after = Pt(10)


# ---------------------------------------------------------------- entry

def build_report(report_type, period_start, period_end, template_path=None,
                 highlights="", team_name="AMS Support Team"):
    """Returns BytesIO of the generated .pptx."""
    branded = bool(template_path)
    if template_path:
        prs = Presentation(template_path)
        # remove existing slides, keep master/layouts (client branding)
        xml_slides = prs.slides._sldIdLst
        for sld in list(xml_slides):
            prs.part.drop_rel(sld.rId)
            xml_slides.remove(sld)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    kind_labels = {
        "DSR": "Daily Status Report",
        "WSR": "Weekly Status Report",
        "MSR": "Monthly Status Report",
    }
    if period_start == period_end:
        period_label = period_start.strftime("%A, %d %B %Y")
    else:
        period_label = (f"{period_start.strftime('%d %b %Y')} – "
                        f"{period_end.strftime('%d %b %Y')}")

    m = compute_metrics(period_start, period_end)

    _title_slide(prs, kind_labels[report_type], period_label, team_name, branded)
    _summary_slide(prs, m, period_label, branded)
    _ageing_slide(prs, m, period_label, branded)
    _resolved_slide(prs, m, period_label, branded)
    _highlights_slide(prs, highlights, period_label, branded)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
